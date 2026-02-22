#!/usr/bin/env python3
"""
Create a PowerPoint presentation.
Extend this for event-day song slides (e.g. one slide per song).
Supports adding audio per slide and advancing after the sound length.
"""

import hashlib
import io
import re
from pathlib import Path

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Inches, Pt

try:
    from mutagen import File as MutagenFile
except ImportError:
    MutagenFile = None

SOURCES_DIR = Path(__file__).resolve().parent / "sources"
SONGS_XLSX = SOURCES_DIR / "music_bingo_songs.xlsx"
BACKGROUND_IMAGE = SOURCES_DIR / "background.jpeg"
# Default folder for audio files (same naming as download_songs: <sanitized_first_column>.m4a).
AUDIO_DIR_DEFAULT = Path(__file__).resolve().parent / "output" / "audio"

# Characters unsafe in filenames (must match download_songs.py).
_UNSAFE_FILENAME = re.compile(r'[<>:"/\\|?*\x00-\x1f]')

# Content box: song/artist/year text is constrained to this rectangle (in inches).
# Adjust to match the area you want (e.g. between logo, mics, and sponsors panel).
# Slide size is 13.333" x 7.5". Measure your drawn square and set:
TEXT_BOX_LEFT = 2.2
TEXT_BOX_TOP = 3.0
TEXT_BOX_WIDTH = 8.9
TEXT_BOX_HEIGHT = 3.45

# Set to True to draw a visible border around the content box (for debugging).
SHOW_TEXT_BOX_BORDER = False


def _sanitize_filename(name: str) -> str:
    """Match download_songs.py: safe filename from first column."""
    return _UNSAFE_FILENAME.sub("_", name).strip() or "unnamed"


def _parse_song_cell(cell_value: str) -> tuple[str, str, str]:
    """
    Parse "SONG_NAME - ARTIST_NAME (RELEASE_YEAR)".
    First " - " is the separator. Returns (song_name, artist_name, release_year).
    """
    if not cell_value or not str(cell_value).strip():
        return ("", "", "")
    s = str(cell_value).strip()
    # First " - " separates song from "ARTIST (YEAR)"
    idx = s.find(" - ")
    if idx < 0:
        return (s, "", "")
    song_name = s[:idx].strip()
    rest = s[idx + 3 :].strip()  # after " - "
    # (RELEASE_YEAR) at the end
    match = re.search(r"\s*\(([^)]+)\)\s*$", rest)
    if match:
        release_year = match.group(1).strip()
        artist_name = rest[: match.start()].strip()
    else:
        release_year = ""
        artist_name = rest
    return (song_name, artist_name, release_year)


def load_songs_from_xlsx(sheet_id: int, xlsx_path: Path = SONGS_XLSX) -> list[tuple[str, str, str, str]]:
    """Load first column of the sheet at index sheet_id (0-based); skip header row.
    Returns list of (song_name, artist_name, release_year, raw_cell_text) for matching audio files."""
    if not xlsx_path.is_file():
        raise FileNotFoundError(f"Songs file not found: {xlsx_path}")
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
    if sheet_id < 0 or sheet_id >= len(wb.sheetnames):
        wb.close()
        names = ", ".join(f"{i}: {n!r}" for i, n in enumerate(wb.sheetnames))
        raise ValueError(f"Sheet id {sheet_id} out of range. Available: {names}")
    ws = wb.worksheets[sheet_id]
    rows = list(ws.iter_rows(min_col=1, max_col=1, values_only=True))
    wb.close()
    result = []
    for i, (cell,) in enumerate(rows):
        if i == 0:
            continue  # skip header row
        if cell is None:
            continue
        raw = str(cell).strip()
        parsed = _parse_song_cell(raw)
        if not parsed[0]:  # skip empty song name
            continue
        result.append((*parsed, raw))
    return result


def _add_background_to_slide(slide, prs, image_path: Path) -> None:
    """Add a full-slide background image and send it to the back."""
    if not image_path.is_file():
        raise FileNotFoundError(f"Background image not found: {image_path}")
    pic = slide.shapes.add_picture(
        str(image_path),
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    pic_element = pic._element
    slide.shapes._spTree.remove(pic_element)
    slide.shapes._spTree.insert(0, pic_element)


class _AudioMedia:
    """Minimal media object for embedding audio (like pptx.media.Video for video)."""

    def __init__(self, blob: bytes, content_type: str, filename: str):
        self._blob = blob
        self._content_type = content_type
        self._filename = filename

    @property
    def blob(self) -> bytes:
        return self._blob

    @property
    def content_type(self) -> str:
        return self._content_type

    @property
    def ext(self) -> str:
        return "m4a" if "m4a" in (self._content_type or "") or (self._filename or "").endswith(".m4a") else "m4a"

    @property
    def sha1(self) -> str:
        """SHA1 hash of the blob (used by package to deduplicate media)."""
        return hashlib.sha1(self._blob).hexdigest()

    @classmethod
    def from_path(cls, path: Path) -> "_AudioMedia":
        with open(path, "rb") as f:
            blob = f.read()
        return cls(blob, "audio/mp4", path.name)


def _get_audio_duration_seconds(audio_path: Path) -> float:
    """Return duration in seconds. Requires mutagen."""
    if not MutagenFile:
        return 0.0
    try:
        af = MutagenFile(str(audio_path))
        return float(af.info.length) if af and af.info else 0.0
    except Exception:
        return 0.0


def _add_audio_to_slide(slide, audio_path: Path) -> None:
    """Embed audio on the slide (speaker icon, plays on show). Uses OOXML audioFile."""
    from pptx.media import SPEAKER_IMAGE_BYTES

    slide_part = slide.part
    package = slide_part.package
    audio = _AudioMedia.from_path(audio_path)
    media_part = package.get_or_add_media_part(audio)
    audio_rId = slide_part.relate_to(media_part, RT.AUDIO)
    _, poster_rId = slide_part.get_or_add_image_part(io.BytesIO(SPEAKER_IMAGE_BYTES))
    shape_id = slide.shapes._spTree.max_shape_id + 1
    name = audio_path.stem or "Audio"
    # Small icon off bottom-right (EMU: 1 inch = 914400)
    x, y = int(12.333 * 914400), int(6.5 * 914400)
    cx, cy = int(0.3 * 914400), int(0.3 * 914400)
    # p:pic with a:audioFile and blipFill (speaker image)
    audio_pic_xml = (
        "<p:pic %s>\n"
        "  <p:nvPicPr>\n"
        '    <p:cNvPr id="%d" name="%s">\n'
        '      <a:hlinkClick r:id="" action="ppaction://media"/>\n'
        "    </p:cNvPr>\n"
        '    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>\n'
        "    <p:nvPr>\n"
        '      <a:audioFile r:link="%s"/>\n'
        "    </p:nvPr>\n"
        "  </p:nvPicPr>\n"
        "  <p:blipFill>\n"
        '    <a:blip r:embed="%s"/>\n'
        "    <a:stretch><a:fillRect/></a:stretch>\n"
        "  </p:blipFill>\n"
        "  <p:spPr>\n"
        "    <a:xfrm>\n"
        '      <a:off x="%d" y="%d"/>\n'
        '      <a:ext cx="%d" cy="%d"/>\n'
        "    </a:xfrm>\n"
        '    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>\n'
        "  </p:spPr>\n"
        "</p:pic>"
        % (nsdecls("a", "p", "r"), shape_id, name.replace("&", "&amp;"), audio_rId, poster_rId, x, y, cx, cy)
    )
    pic_el = parse_xml(audio_pic_xml)
    slide.shapes._spTree.append(pic_el)


def _set_slide_advance_after(slide, advance_milliseconds: int) -> None:
    """Set slide to advance automatically after the given number of milliseconds."""
    sld = slide._element
    trans_xml = (
        '<p:transition %s advTm="%d" advClick="1">\n'
        "  <p:fade/>\n"
        "</p:transition>"
        % (nsdecls("p"), advance_milliseconds)
    )
    trans_el = parse_xml(trans_xml)
    for i, child in enumerate(sld):
        if child.tag.endswith("}transition"):
            sld.replace(child, trans_el)
            return
    # Insert before timing (p:timing) or at 2 (after cSld, clrMapOvr)
    insert_at = 2
    for i, child in enumerate(sld):
        if child.tag.endswith("}timing"):
            insert_at = i
            break
        if child.tag.endswith("}cSld"):
            insert_at = i + 1
        elif child.tag.endswith("}clrMapOvr"):
            insert_at = i + 1
    sld.insert(insert_at, trans_el)


def create_sample_presentation(output_path: str | Path = "output.pptx") -> Path:
    """Create a sample presentation with title and content slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide (layout 0)
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    _add_background_to_slide(slide, prs, BACKGROUND_IMAGE)
    slide.shapes.title.text = "Music Bingo"
    slide.placeholders[1].text = "Event day"

    # Blank layout for custom content (layout 6 is often blank)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _add_background_to_slide(slide, prs, BACKGROUND_IMAGE)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sample content slide"
    p.font.size = Pt(44)
    p.font.bold = True

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def create_slides_from_songs(
    song_titles: list[str],
    output_path: str | Path = "songs_slides.pptx",
    title: str = "Songs",
) -> Path:
    """
    Create one slide per song (e.g. for event-day display).
    song_titles: list of song names to show, one per slide.
    """
    # Convert to (song, artist, year) with empty artist/year for backward compat
    triples = [(t.strip(), "", "") for t in song_titles]
    return create_slides_from_song_triples(triples, output_path)


def create_slides_from_song_triples(
    songs: list[tuple[str, str, str] | tuple[str, str, str, str]],
    output_path: str | Path = "songs_slides.pptx",
    audio_dir: Path | None = None,
) -> Path:
    """
    Create one slide per song. Each slide: SONG NAME (big), ARTIST (smaller), YEAR (under artist).
    songs: list of (song_name, artist_name, release_year) or (song_name, artist_name, release_year, raw_cell).
    If audio_dir is set and each item has raw_cell, looks for <sanitized raw_cell>.m4a in audio_dir,
    embeds it on the slide and sets advance-after to the sound length.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    audio_dir_path = Path(audio_dir) if audio_dir else None

    # Intro slide at the start: "ROUND" / "Starting soon..."
    intro = prs.slides.add_slide(blank)
    _add_background_to_slide(intro, prs, BACKGROUND_IMAGE)
    left = Inches(TEXT_BOX_LEFT)
    top = Inches(TEXT_BOX_TOP)
    width = Inches(TEXT_BOX_WIDTH)
    height = Inches(TEXT_BOX_HEIGHT)
    intro_box = intro.shapes.add_textbox(left, top, width, height)
    intro_tf = intro_box.text_frame
    intro_tf.word_wrap = True
    intro_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    intro_p0 = intro_tf.paragraphs[0]
    intro_p0.text = "ROUND"
    intro_p0.font.size = Pt(64)
    intro_p0.font.bold = True
    intro_p0.alignment = PP_ALIGN.CENTER
    intro_p1 = intro_tf.add_paragraph()
    intro_p1.text = "Starting soon..."
    intro_p1.font.size = Pt(38)
    intro_p1.alignment = PP_ALIGN.CENTER
    intro_p1.space_before = Pt(12)

    for i, row in enumerate(songs, start=1):
        if len(row) == 4:
            song_name, artist_name, release_year, raw_cell = row
        else:
            song_name, artist_name, release_year = row[0], row[1], row[2]
            raw_cell = ""
        slide = prs.slides.add_slide(blank)
        _add_background_to_slide(slide, prs, BACKGROUND_IMAGE)

        if audio_dir_path and raw_cell:
            safe_name = _sanitize_filename(raw_cell)
            audio_path = audio_dir_path / f"{safe_name}.m4a"
            if audio_path.is_file():
                _add_audio_to_slide(slide, audio_path)
                duration_sec = _get_audio_duration_seconds(audio_path)
                if duration_sec > 0:
                    _set_slide_advance_after(slide, int(round(duration_sec * 1000)))

        # Text fits only inside the configured content box (TEXT_BOX_*).
        left = Inches(TEXT_BOX_LEFT)
        top = Inches(TEXT_BOX_TOP)
        width = Inches(TEXT_BOX_WIDTH)
        height = Inches(TEXT_BOX_HEIGHT)

        if SHOW_TEXT_BOX_BORDER:
            # Debug: draw a visible rectangle so you can see the content box.
            border_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            border_shape.fill.background()
            border_shape.line.color.rgb = RGBColor(255, 0, 0)
            border_shape.line.width = Pt(2)
            # Send border behind the text box (insert at 0 after background).
            border_el = border_shape._element
            slide.shapes._spTree.remove(border_el)
            slide.shapes._spTree.insert(1, border_el)  # after background pic (0)

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink text to fit the box

        # Song name – large but sized to fit zone
        p0 = tf.paragraphs[0]
        p0.text = song_name
        p0.font.size = Pt(64)
        p0.font.bold = True
        p0.alignment = PP_ALIGN.CENTER

        if artist_name or release_year:
            # Artist name – under song
            p1 = tf.add_paragraph()
            p1.text = artist_name
            p1.font.size = Pt(38)
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(8)

        if release_year:
            # Release year – under artist
            p2 = tf.add_paragraph()
            p2.text = release_year
            p2.font.size = Pt(30)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(5)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    import argparse
    import sys

    out_dir = Path(__file__).resolve().parent / "output"
    out_dir.mkdir(exist_ok=True)

    parser = argparse.ArgumentParser(description="Create Music Bingo PowerPoint slides.")
    parser.add_argument(
        "-o", "--output",
        metavar="NAME",
        help="Output file path (default: output/sample.pptx or output/songs_slides.pptx)",
    )
    parser.add_argument(
        "-s", "--sheet",
        type=int,
        metavar="ID",
        help="Sheet index (0-based): read songs from first column of sources/music_bingo_songs.xlsx.",
    )
    parser.add_argument(
        "-a", "--audio-dir",
        metavar="DIR",
        help="Directory containing .m4a files named like the first column (e.g. output from download_songs.py). Adds audio to each slide and sets advance-after to sound length.",
    )
    parser.add_argument(
        "songs",
        nargs="*",
        help="Song list: file path (one song per line) or song names as arguments (ignored if -s used)",
    )
    args = parser.parse_args()

    if args.output:
        p = Path(args.output)
        output_path = (out_dir / p.name) if p.parent == Path(".") else p
        if output_path.suffix.lower() != ".pptx":
            output_path = output_path.with_suffix(".pptx")
    else:
        output_path = out_dir / "sample.pptx"

    if args.sheet is not None:
        try:
            songs = load_songs_from_xlsx(args.sheet)
        except ValueError as e:
            sys.exit(str(e))
        if not songs:
            sys.exit(f"No songs found in sheet {args.sheet}")
        if not args.output:
            output_path = out_dir / "songs_slides.pptx"
        audio_dir = Path(args.audio_dir) if args.audio_dir else None
        create_slides_from_song_triples(songs, output_path, audio_dir=audio_dir)
        print(f"Saved song slides ({len(songs)} slides): {output_path}")
    elif args.songs:
        if len(args.songs) == 1 and Path(args.songs[0]).is_file():
            song_titles = Path(args.songs[0]).read_text(encoding="utf-8").strip().splitlines()
        else:
            song_titles = args.songs
        if not args.output:
            output_path = out_dir / "songs_slides.pptx"
        create_slides_from_songs(song_titles, output_path)
        print(f"Saved song slides: {output_path}")
    else:
        if not args.output:
            output_path = out_dir / "sample.pptx"
        create_sample_presentation(output_path)
        print(f"Saved: {output_path}")
